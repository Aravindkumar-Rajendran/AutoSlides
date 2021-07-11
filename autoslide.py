import os
import sys
import random
import torch
import torchvision
from torchvision.models.detection.faster_rcnn import FastRCNNPredictor
from torchvision.models.detection.mask_rcnn import MaskRCNNPredictor
from torchvision.transforms import transforms
import argparse
import easyocr
reader = easyocr.Reader(['en'])

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('template.pptx')
blank_slide_layout = prs.slide_layouts[1]



import cv2
import numpy as np

from utils import (
    overlay_ann,
    overlay_mask,
    show
)

seed = 1234
random.seed(seed)
torch.manual_seed(seed)
torch.cuda.manual_seed_all(seed)
torch.backends.cudnn.deterministic = True
torch.backends.cudnn.benchmark = False



CATEGORIES2LABELS = {
    0: "bg",
    1: "text",
    2: "title",
    3: "list",
    4: "table",
    5: "figure"
}
SAVE_PATH = "output/"
MODEL_PATH = "models/model_196000.pth"
OUTPUT_PATH = "output"

def get_instance_segmentation_model(num_classes):
    model = torchvision.models.detection.maskrcnn_resnet50_fpn(pretrained=True)
    in_features = model.roi_heads.box_predictor.cls_score.in_features
    model.roi_heads.box_predictor = FastRCNNPredictor(in_features, num_classes)
    in_features_mask = model.roi_heads.mask_predictor.conv5_mask.in_channels
    hidden_layer = 256

    model.roi_heads.mask_predictor = MaskRCNNPredictor(
        in_features_mask,
        hidden_layer,
        num_classes
    )
    return model


def create_pptx(images):
    num_classes = 6
    model = get_instance_segmentation_model(num_classes)
    model.cuda()

    if os.path.exists(MODEL_PATH):
        checkpoint_path = MODEL_PATH

    print(checkpoint_path)
    assert os.path.exists(checkpoint_path)
    checkpoint = torch.load(checkpoint_path, map_location='cpu')
    model.load_state_dict(checkpoint['model'])
    model.eval()

    if os.path.exists(images[0]):
        image_path = images[0]

    print(image_path)
    assert os.path.exists(image_path)

    image_raw = cv2.imread(image_path)
    image = image_raw.copy()
    print("Image: ", image.shape)
    rat = 1000 / image.shape[0]
    image = cv2.resize(image, None, fx=rat, fy=rat)

    transform = transforms.Compose([
        transforms.ToPILImage(),
        transforms.ToTensor()
    ])
    image = transform(image)

    with torch.no_grad():
        prediction = model([image.cuda()])

    image = torch.squeeze(image, 0).permute(1, 2, 0).mul(255).numpy().astype(np.uint8)

    results = []
    titles = []
    texts = []
    with open('OCR.txt', "a+") as file:
        for pred in prediction:
            for idx, mask in enumerate(pred['masks']):
                if pred['scores'][idx].item() < 0.7:
                    continue

                m = mask[0].mul(255).byte().cpu().numpy()
                box = list(map(int, pred["boxes"][idx].tolist()))
                label = CATEGORIES2LABELS[pred["labels"][idx].item()]

                score = pred["scores"][idx].item()

                # image = overlay_mask(image, m)
                image = overlay_ann(image, m, box, label, score)
                # print(box, label, score, image.shape)
                
                part = image_raw[int(box[1]/rat):int(box[3]/rat), int(box[0]/rat):int(box[2]/rat)]
                # print(int(box[1]/rat), int(box[3]/rat), int(box[0]/rat), int(box[2]/rat))
                file_name = f"output/{str(idx)}.jpg"
                cv2.imwrite(file_name, part)
                if label in ['text', 'title', 'list']:  # List to be added
                    result = reader.readtext(part, paragraph=True)
                    # print(result)
                    file.write(str(idx) + " " + label)
                    file.write("\n")
                    file.writelines(result[0][1])
                    file.write("\n")
                    results.append((box, score, label, result[0][1]))
                    if label == 'title':
                        titles.append(''.join([i for i in result[0][1] if not (i.isdigit() or i in [',', '.'])]))
                    else:
                        texts.append(''.join([i for i in result[0][1] if not (i.isdigit() or i in [',', '.'])]))
                else:
                    results.append((box, score, label, file_name))

    sorted_results = sorted(results, key=lambda x:x[0][1])

    column_sorted_results = []
    second_column = []
    for result in sorted_results:
        if result[0][0] <= 350:
            column_sorted_results.append(result)
        else:
            second_column.append(result)

    column_sorted_results += second_column

    same_slide = True
    previous_label = ''
    print("results : ", column_sorted_results)

    for box, score, label, item in column_sorted_results:
        if label== "title":
            slide = prs.slides.add_slide(blank_slide_layout)
            top = Inches(1)
            left = Inches(1)
            width = height = Inches(7)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(28)
            same_slide = True
            
        elif label=="figure":
            if not same_slide:
                slide = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(1)
            height = Inches(4) 
            pic = slide.shapes.add_picture(item, left, top, height=height)
            same_slide = True

        elif label in ['text', 'list']:
            if previous_label == "figure" and label == "text":
                top = Inches(6)
                left = Inches(2)
                width = height = Inches(5)

            else:
                if not (len(item) > 50) and not same_slide:
                    top = Inches(1)
                    left = Inches(1)
                    width = height = Inches(7)
                else:
                    top = Inches(2)
                    left = Inches(1)
                    width = height = Inches(7)

            if not same_slide:
                slide = prs.slides.add_slide(blank_slide_layout)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)

            same_slide = False if (len(item) > 50 or previous_label == 'figure') else True #second title

        previous_label = label


    if os.path.exists(OUTPUT_PATH):
        cv2.imwrite(OUTPUT_PATH+'/{}'.format(os.path.basename(image_path)), image)
    else:
        os.mkdir(OUTPUT_PATH)
        cv2.imwrite(OUTPUT_PATH+'/{}'.format(os.path.basename(image_path)), image)

    # show(image)
    file_name = os.path.join(OUTPUT_PATH, str(random.randint(1111, 9999)) +  '.pptx' )
    prs.save(file_name)

    output_images = OUTPUT_PATH+'/{}'.format(os.path.basename(image_path))

    return output_images, file_name, titles, texts



if __name__ == "__main__":
    create_pptx()
